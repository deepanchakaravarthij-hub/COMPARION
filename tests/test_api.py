from __future__ import annotations

import os
from dataclasses import replace
from io import BytesIO
from pathlib import Path
from unittest.mock import patch
from uuid import uuid4

import fitz  # type: ignore[import-untyped]
from docx import Document
from fastapi.testclient import TestClient
from openpyxl import Workbook  # type: ignore[import-untyped]
from openpyxl.styles import Font  # type: ignore[import-untyped]
from PIL import Image, ImageDraw
from pptx import Presentation  # type: ignore[import-untyped]
from pptx.enum.shapes import MSO_SHAPE  # type: ignore[import-untyped]
from pptx.util import Inches, Pt  # type: ignore[import-untyped]

from app.core.config import get_settings
from app.main import app
from app.services.compare_service import compare_files
from app.services.job_store import add_audit_event

client = TestClient(app)


def _image_bytes(color: str = "white", mark: bool = False) -> bytes:
    image = Image.new("RGB", (100, 100), color)
    if mark:
        draw = ImageDraw.Draw(image)
        draw.rectangle((20, 20, 50, 50), fill="black")
    output = BytesIO()
    image.save(output, format="PNG")
    return output.getvalue()


def _pdf_bytes(text: str) -> bytes:
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), text)
    content = bytes(document.tobytes())
    document.close()
    return content


def _docx_bytes(
    paragraph_runs: list[list[tuple[str, bool]]],
    table_rows: list[list[str]] | None = None,
) -> bytes:
    document = Document()
    for runs in paragraph_runs:
        paragraph = document.add_paragraph()
        for text, bold in runs:
            run = paragraph.add_run(text)
            run.bold = bold

    if table_rows:
        table = document.add_table(rows=len(table_rows), cols=len(table_rows[0]))
        for row_index, row in enumerate(table_rows):
            for column_index, value in enumerate(row):
                table.cell(row_index, column_index).text = value

    output = BytesIO()
    document.save(output)
    return output.getvalue()


def _xlsx_bytes(
    *,
    first_sheet_name: str = "Summary",
    first_sheet_order: int = 0,
    second_sheet_name: str = "Data",
    second_sheet_order: int = 1,
    amount: int = 100,
    formula: str = "=B2*2",
    hidden_row: bool = False,
    hidden_column: bool = False,
    hidden_sheet: bool = False,
    bold_amount: bool = False,
) -> bytes:
    workbook = Workbook()
    first_sheet = workbook.active
    first_sheet.title = first_sheet_name
    second_sheet = workbook.create_sheet(second_sheet_name)
    sheets = [first_sheet, second_sheet]
    workbook._sheets = [sheets[first_sheet_order], sheets[second_sheet_order]]

    first_sheet["A1"] = "metric"
    first_sheet["B1"] = "value"
    first_sheet["A2"] = "amount"
    first_sheet["B2"] = amount
    first_sheet["C2"] = formula
    first_sheet["B2"].number_format = "#,##0"
    if bold_amount:
        first_sheet["B2"].font = Font(bold=True)

    second_sheet["A1"] = "status"
    second_sheet["B1"] = "ok"

    if hidden_row:
        first_sheet.row_dimensions[2].hidden = True
    if hidden_column:
        first_sheet.column_dimensions["B"].hidden = True
    if hidden_sheet:
        second_sheet.sheet_state = "hidden"

    output = BytesIO()
    workbook.save(output)
    return output.getvalue()


def _pptx_bytes(
    *,
    first_title: str = "Overview",
    first_text: str = "Revenue increased",
    first_bold: bool = False,
    first_table_value: str = "100",
    second_title: str = "Details",
    second_text: str = "Risk unchanged",
    second_bold: bool = False,
    second_table_value: str = "Stable",
    image_color: str = "black",
    slide_order: tuple[int, int] = (0, 1),
) -> bytes:
    presentation = Presentation()

    slide_one = presentation.slides.add_slide(presentation.slide_layouts[5])
    title_one = slide_one.shapes.title
    if title_one is not None:
        title_one.text = first_title
    textbox_one = slide_one.shapes.add_textbox(Inches(1), Inches(1.5), Inches(6), Inches(1))
    run_one = textbox_one.text_frame.paragraphs[0].add_run()
    run_one.text = first_text
    run_one.font.bold = first_bold
    run_one.font.size = Pt(20)
    table_one = slide_one.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(2.5),
        Inches(5),
        Inches(1.5),
    ).table
    table_one.cell(0, 0).text = "Metric"
    table_one.cell(0, 1).text = "Value"
    table_one.cell(1, 0).text = "Score"
    table_one.cell(1, 1).text = first_table_value
    slide_one.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6),
        Inches(1.5),
        Inches(2),
        Inches(1),
    )
    slide_one.shapes.add_picture(BytesIO(_image_bytes(color=image_color)), Inches(6), Inches(3))

    slide_two = presentation.slides.add_slide(presentation.slide_layouts[5])
    title_two = slide_two.shapes.title
    if title_two is not None:
        title_two.text = second_title
    textbox_two = slide_two.shapes.add_textbox(Inches(1), Inches(1.5), Inches(6), Inches(1))
    run_two = textbox_two.text_frame.paragraphs[0].add_run()
    run_two.text = second_text
    run_two.font.bold = second_bold
    table_two = slide_two.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(2.5),
        Inches(5),
        Inches(1.5),
    ).table
    table_two.cell(0, 0).text = "Label"
    table_two.cell(0, 1).text = "Status"
    table_two.cell(1, 0).text = "Risk"
    table_two.cell(1, 1).text = second_table_value

    _reorder_pptx_slides(presentation, slide_order)

    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _reorder_pptx_slides(presentation: Presentation, order: tuple[int, int]) -> None:
    slide_ids = list(presentation.slides._sldIdLst)
    reordered = [slide_ids[index] for index in order]
    presentation.slides._sldIdLst[:] = reordered


def test_health() -> None:
    res = client.get("/health")
    assert res.status_code == 200
    assert res.json()["status"] == "ok"
    assert res.json()["environment"] == "local"
    assert "X-Request-ID" in res.headers


def test_viewer_page_is_served() -> None:
    res = client.get("/viewer")
    assert res.status_code == 200
    assert "COMPARION Viewer" in res.text


def test_request_id_is_preserved() -> None:
    res = client.get("/health", headers={"X-Request-ID": "test-request-id"})
    assert res.status_code == 200
    assert res.headers["X-Request-ID"] == "test-request-id"


def test_compare_pdf_and_fetch_result_and_report() -> None:
    files = {
        "file_a": ("a.pdf", _pdf_bytes("hello contract"), "application/pdf"),
        "file_b": ("b.pdf", _pdf_bytes("hello updated contract"), "application/pdf"),
    }
    compare_res = client.post("/v1/compare", files=files)
    assert compare_res.status_code == 200
    assert compare_res.json()["request_id"]

    job_id = compare_res.json()["job_id"]
    status_res = client.get(f"/v1/jobs/{job_id}")
    assert status_res.status_code == 200
    assert status_res.json()["status"] == "completed"
    assert status_res.json()["file_a_type"] == "pdf"

    result_res = client.get(f"/v1/jobs/{job_id}/result")
    assert result_res.status_code == 200
    payload = result_res.json()
    assert payload["result_schema_version"] == "2.1"
    assert payload["file_type"] == "pdf"
    assert payload["changes"]

    report_res = client.get(f"/v1/jobs/{job_id}/report.html")
    assert report_res.status_code == 200
    assert "COMPARION Report" in report_res.text


def test_compare_docx_and_fetch_result_and_report() -> None:
    file_a = _docx_bytes(
        [
            [("Agreement summary", False)],
            [("Payment due in 30 days", False)],
            [("Reviewed by Legal", False)],
        ],
        [["Clause", "Value"], ["Fee", "100"]],
    )
    file_b = _docx_bytes(
        [
            [("Agreement summary", False)],
            [("Payment due in 45 days", False)],
            [("Reviewed by Legal", True)],
        ],
        [["Clause", "Value"], ["Fee", "125"]],
    )
    files = {
        "file_a": (
            "a.docx",
            file_a,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ),
        "file_b": (
            "b.docx",
            file_b,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ),
    }

    compare_res = client.post("/v1/compare", files=files)
    assert compare_res.status_code == 200

    job_id = compare_res.json()["job_id"]
    status_res = client.get(f"/v1/jobs/{job_id}")
    assert status_res.status_code == 200
    assert status_res.json()["status"] == "completed"
    assert status_res.json()["file_a_type"] == "docx"

    result_res = client.get(f"/v1/jobs/{job_id}/result")
    assert result_res.status_code == 200
    payload = result_res.json()
    categories = {change["category"] for change in payload["changes"]}
    assert {"text", "formatting", "table"}.issubset(categories)
    assert payload["diagnostics"]["docx"]["paragraph_count"]["a"] >= 3

    report_res = client.get(f"/v1/jobs/{job_id}/report.html")
    assert report_res.status_code == 200
    assert "Text changes" in report_res.text
    assert "Formatting changes" in report_res.text
    assert "Table changes" in report_res.text


def test_compare_xlsx_and_fetch_result_and_report() -> None:
    file_a = _xlsx_bytes(
        first_sheet_name="Summary",
        first_sheet_order=0,
        second_sheet_name="Data",
        second_sheet_order=1,
        amount=100,
        formula="=B2*2",
    )
    file_b = _xlsx_bytes(
        first_sheet_name="Summary",
        first_sheet_order=1,
        second_sheet_name="DataRenamed",
        second_sheet_order=0,
        amount=125,
        formula="=B2*3",
        hidden_row=True,
        hidden_column=True,
        hidden_sheet=True,
        bold_amount=True,
    )
    files = {
        "file_a": (
            "a.xlsx",
            file_a,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ),
        "file_b": (
            "b.xlsx",
            file_b,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ),
    }

    compare_res = client.post("/v1/compare", files=files)
    assert compare_res.status_code == 200

    job_id = compare_res.json()["job_id"]
    status_res = client.get(f"/v1/jobs/{job_id}")
    assert status_res.status_code == 200
    assert status_res.json()["status"] == "completed"
    assert status_res.json()["file_a_type"] == "xlsx"

    result_res = client.get(f"/v1/jobs/{job_id}/result")
    assert result_res.status_code == 200
    payload = result_res.json()
    categories = {change["category"] for change in payload["changes"]}
    assert {"sheet", "formula", "text", "structure", "metadata"}.issubset(categories)
    assert payload["diagnostics"]["xlsx"]["changed_cells"]

    report_res = client.get(f"/v1/jobs/{job_id}/report.html")
    assert report_res.status_code == 200
    assert "Sheet summary" in report_res.text
    assert "Changed cells" in report_res.text


def test_compare_pptx_and_fetch_result_and_report() -> None:
    file_a = _pptx_bytes(
        first_title="Overview",
        first_text="Revenue update",
        first_bold=False,
        first_table_value="100",
        second_title="Details",
        second_text="Risk unchanged",
        second_bold=False,
        second_table_value="Stable",
        image_color="black",
        slide_order=(0, 1),
    )
    file_b = _pptx_bytes(
        first_title="Overview",
        first_text="Revenue update",
        first_bold=True,
        first_table_value="80",
        second_title="Details Updated",
        second_text="Risk elevated",
        second_bold=True,
        second_table_value="High",
        image_color="red",
        slide_order=(1, 0),
    )
    files = {
        "file_a": (
            "a.pptx",
            file_a,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ),
        "file_b": (
            "b.pptx",
            file_b,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ),
    }

    compare_res = client.post("/v1/compare", files=files)
    assert compare_res.status_code == 200

    job_id = compare_res.json()["job_id"]
    status_res = client.get(f"/v1/jobs/{job_id}")
    assert status_res.status_code == 200
    assert status_res.json()["status"] == "completed"
    assert status_res.json()["file_a_type"] == "pptx"

    result_res = client.get(f"/v1/jobs/{job_id}/result")
    assert result_res.status_code == 200
    payload = result_res.json()
    categories = {change["category"] for change in payload["changes"]}
    assert {"text", "formatting", "table", "image", "structure"}.issubset(categories)
    assert payload["diagnostics"]["pptx"]["slide_summary"]
    assert payload["diagnostics"]["pptx"]["changed_objects"]

    report_res = client.get(f"/v1/jobs/{job_id}/report.html")
    assert report_res.status_code == 200
    assert "Slide summary" in report_res.text
    assert "Object changes" in report_res.text


def test_compare_identical_images_has_no_changes() -> None:
    content = _image_bytes()
    files = {
        "file_a": ("a.png", content, "image/png"),
        "file_b": ("b.png", content, "image/png"),
    }
    compare_res = client.post("/v1/compare", files=files)
    assert compare_res.status_code == 200

    job_id = compare_res.json()["job_id"]
    result_res = client.get(f"/v1/jobs/{job_id}/result")
    assert result_res.status_code == 200
    payload = result_res.json()
    assert payload["file_type"] == "image"
    assert payload["changes"] == []
    assert payload["summary"] == "No differences detected"


def test_job_history_and_filtered_result() -> None:
    files = {
        "file_a": ("a.pdf", _pdf_bytes("alpha text"), "application/pdf"),
        "file_b": ("b.pdf", _pdf_bytes("beta text"), "application/pdf"),
    }
    compare_res = client.post("/v1/compare", files=files)
    assert compare_res.status_code == 200
    job_id = compare_res.json()["job_id"]

    history_res = client.get("/v1/jobs?limit=10")
    assert history_res.status_code == 200
    history_payload = history_res.json()
    assert history_payload["items"]
    assert any(item["job_id"] == job_id for item in history_payload["items"])

    result_res = client.get(f"/v1/jobs/{job_id}/result?category=text&limit=5&offset=0")
    assert result_res.status_code == 200
    payload = result_res.json()
    assert payload["pagination"]["limit"] == 5
    assert payload["pagination"]["offset"] == 0
    assert payload["pagination"]["total_filtered"] >= 1
    assert all(change["category"] == "text" for change in payload["changes"])


def test_result_version_compatibility_view() -> None:
    files = {
        "file_a": (
            "a.xlsx",
            _xlsx_bytes(amount=100, formula="=B2*2"),
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ),
        "file_b": (
            "b.xlsx",
            _xlsx_bytes(amount=120, formula="=B2*3"),
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ),
    }
    compare_res = client.post("/v1/compare", files=files)
    assert compare_res.status_code == 200
    job_id = compare_res.json()["job_id"]

    latest_res = client.get(f"/v1/jobs/{job_id}/result?result_schema_version=2.1")
    assert latest_res.status_code == 200
    assert latest_res.json()["result_schema_version"] == "2.1"

    legacy_res = client.get(f"/v1/jobs/{job_id}/result?result_schema_version=2.0")
    assert legacy_res.status_code == 200
    legacy = legacy_res.json()
    assert legacy["result_schema_version"] == "2.0"
    assert "udm" not in legacy or legacy["udm"] is None
    for change in legacy["changes"]:
        assert change["category"] not in {"formula", "sheet", "slide"}
        assert change["source_ref"].get("sheet") is None
        assert change["source_ref"].get("cell") is None
        assert change["source_ref"].get("slide") is None


def test_report_link_unsigned_and_signed_modes() -> None:
    files = {
        "file_a": (
            "a.docx",
            _docx_bytes([[("Alpha", False)]], None),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ),
        "file_b": (
            "b.docx",
            _docx_bytes([[("Beta", False)]], None),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ),
    }
    compare_res = client.post("/v1/compare", files=files)
    assert compare_res.status_code == 200
    job_id = compare_res.json()["job_id"]

    link_res = client.get(f"/v1/jobs/{job_id}/report-link")
    assert link_res.status_code == 200
    link_payload = link_res.json()
    assert link_payload["signed"] is False
    assert link_payload["url"].endswith(f"/v1/jobs/{job_id}/report.html")

    signed_settings = replace(get_settings(), object_storage_enabled=True)
    with patch("app.api.routes.get_settings", return_value=signed_settings):
        link_signed_res = client.get(f"/v1/jobs/{job_id}/report-link")
        assert link_signed_res.status_code == 200
        signed_payload = link_signed_res.json()
        assert signed_payload["signed"] is True
        assert "token=" in signed_payload["url"]


def test_rejects_unsupported_extension() -> None:
    files = {
        "file_a": ("a.txt", b"hello", "text/plain"),
        "file_b": ("b.txt", b"hello", "text/plain"),
    }
    res = client.post("/v1/compare", files=files)
    assert res.status_code == 415


def test_rejects_mismatched_supported_types() -> None:
    files = {
        "file_a": ("a.pdf", _pdf_bytes("hello"), "application/pdf"),
        "file_b": ("b.png", _image_bytes(), "image/png"),
    }
    res = client.post("/v1/compare", files=files)
    assert res.status_code == 400


def test_failed_job_includes_error_reason() -> None:
    files = {
        "file_a": ("a.pdf", b"not a pdf", "application/pdf"),
        "file_b": ("b.pdf", b"also not a pdf", "application/pdf"),
    }
    compare_res = client.post("/v1/compare", files=files)
    assert compare_res.status_code == 200

    job_id = compare_res.json()["job_id"]
    status_res = client.get(f"/v1/jobs/{job_id}")
    assert status_res.status_code == 200
    payload = status_res.json()
    assert payload["status"] == "failed"
    assert payload["error"]["code"] == "internal_error"
    assert payload["error"]["message"]

    result_res = client.get(f"/v1/jobs/{job_id}/result")
    assert result_res.status_code == 409


def test_image_diff_is_deterministic() -> None:
    result_a = compare_files("a.png", "b.png", _image_bytes(), _image_bytes(mark=True))
    result_b = compare_files("a.png", "b.png", _image_bytes(), _image_bytes(mark=True))
    assert result_a["changes"] == result_b["changes"]
    assert result_a["changes"][0]["bbox"] == {
        "page": 1,
        "x": 0.61,
        "y": 0.5,
        "width": 0.32,
        "height": 0.32,
    }


def test_docx_diff_is_deterministic() -> None:
    file_a = _docx_bytes(
        [[("Alpha", False)], [("Beta", False)]],
        [["Key", "Value"], ["Limit", "10"]],
    )
    file_b = _docx_bytes(
        [[("Beta", False)], [("Alpha changed", False)]],
        [["Key", "Value"], ["Limit", "20"]],
    )
    result_a = compare_files("a.docx", "b.docx", file_a, file_b)
    result_b = compare_files("a.docx", "b.docx", file_a, file_b)
    assert result_a["changes"] == result_b["changes"]
    assert {change["category"] for change in result_a["changes"]} >= {"text", "table", "structure"}


def test_xlsx_diff_is_deterministic() -> None:
    file_a = _xlsx_bytes(amount=100, formula="=B2*2")
    file_b = _xlsx_bytes(
        first_sheet_name="Summary",
        first_sheet_order=1,
        second_sheet_name="DataRenamed",
        second_sheet_order=0,
        amount=140,
        formula="=B2*3",
        hidden_row=True,
        hidden_column=True,
        hidden_sheet=True,
    )
    result_a = compare_files("a.xlsx", "b.xlsx", file_a, file_b)
    result_b = compare_files("a.xlsx", "b.xlsx", file_a, file_b)
    assert result_a["changes"] == result_b["changes"]
    assert {change["category"] for change in result_a["changes"]} >= {
        "sheet",
        "formula",
        "text",
        "structure",
    }


def test_pptx_diff_is_deterministic() -> None:
    file_a = _pptx_bytes(
        first_title="Overview",
        first_text="Revenue update",
        first_table_value="100",
        second_title="Details",
        second_text="Risk unchanged",
        second_table_value="Stable",
        image_color="black",
        slide_order=(0, 1),
    )
    file_b = _pptx_bytes(
        first_title="Overview",
        first_text="Revenue update",
        first_bold=True,
        first_table_value="80",
        second_title="Details Updated",
        second_text="Risk elevated",
        second_bold=True,
        second_table_value="High",
        image_color="red",
        slide_order=(1, 0),
    )
    result_a = compare_files("a.pptx", "b.pptx", file_a, file_b)
    result_b = compare_files("a.pptx", "b.pptx", file_a, file_b)
    assert result_a["changes"] == result_b["changes"]
    assert {change["category"] for change in result_a["changes"]} >= {
        "text",
        "formatting",
        "table",
        "image",
        "structure",
    }


def test_semantic_layer_includes_labels_matches_and_risk_summary() -> None:
    file_a = _docx_bytes(
        [
            [("Payment amount is 100 USD", False)],
            [("Liability clause remains unchanged", False)],
            [("Policy audit date is 2024-01-01", False)],
        ],
        [["Clause", "Value"], ["Fee", "100"]],
    )
    file_b = _docx_bytes(
        [
            [("Payment amount is 120 USD", False)],
            [("Liability clause is revised for termination", False)],
            [("Policy audit date is 2024-03-01", False)],
        ],
        [["Clause", "Value"], ["Fee", "120"]],
    )
    result = compare_files("a.docx", "b.docx", file_a, file_b)
    assert result["semantic"]["summary"]
    assert result["semantic"]["provenance"]["method"] == "local-semantic-v1"
    assert result["semantic"]["prompt_templates"]["change_summary"]
    assert result["semantic"]["risk_rule_templates"]["finance"]["keywords"]
    assert "high_risk_count" in result["semantic"]["risk_summary"]
    assert result["semantic"]["risk_summary"]["high_risk_count"] >= 1
    assert any(change.get("semantic_label") for change in result["changes"])
    assert all("semantic_score" in change for change in result["changes"])


def test_semantic_metadata_not_returned_in_legacy_view() -> None:
    files = {
        "file_a": ("a.pdf", _pdf_bytes("Agreement payment fee 100"), "application/pdf"),
        "file_b": ("b.pdf", _pdf_bytes("Agreement payment fee 120"), "application/pdf"),
    }
    compare_res = client.post("/v1/compare", files=files)
    assert compare_res.status_code == 200
    job_id = compare_res.json()["job_id"]

    legacy_res = client.get(f"/v1/jobs/{job_id}/result?result_schema_version=2.0")
    assert legacy_res.status_code == 200
    payload = legacy_res.json()
    assert payload["result_schema_version"] == "2.0"
    assert payload.get("semantic") is None
    for change in payload["changes"]:
        assert "semantic_label" not in change
        assert "semantic_score" not in change


def test_compare_supports_idempotency_key_reuse() -> None:
    key = f"req-{uuid4().hex}"
    files = {
        "file_a": ("a.pdf", _pdf_bytes("alpha"), "application/pdf"),
        "file_b": ("b.pdf", _pdf_bytes("beta"), "application/pdf"),
    }
    first = client.post(f"/v1/compare?idempotency_key={key}", files=files)
    assert first.status_code == 200
    second = client.post(f"/v1/compare?idempotency_key={key}", files=files)
    assert second.status_code == 200
    assert first.json()["job_id"] == second.json()["job_id"]


def test_auth_required_blocks_unauthorized_requests() -> None:
    auth_settings = replace(get_settings(), auth_required=True, auth_api_key="secret-key")
    with patch("app.api.routes.get_settings", return_value=auth_settings):
        blocked = client.get("/v1/jobs")
        assert blocked.status_code == 401
        allowed = client.get("/v1/jobs", headers={"X-API-Key": "secret-key"})
        assert allowed.status_code == 200


def test_metrics_endpoint_returns_alert_shape() -> None:
    metrics = client.get("/v1/metrics")
    assert metrics.status_code == 200
    payload = metrics.json()
    assert "job_status_counts" in payload
    assert "queue_depth" in payload
    assert "failure_rate" in payload
    assert "alerts" in payload


def test_retention_cleanup_endpoint() -> None:
    jobs_root = Path(get_settings().storage_root) / "jobs"
    old_job_dir = jobs_root / "old-test-job"
    old_job_dir.mkdir(parents=True, exist_ok=True)
    old_file = old_job_dir / "artifact.txt"
    old_file.write_text("expired", encoding="utf-8")

    os.utime(old_job_dir, (0, 0))
    os.utime(old_file, (0, 0))
    add_audit_event(event_type="test_setup", details={"name": "retention_cleanup"})

    cleanup_settings = replace(get_settings(), retention_days=0)
    with patch("app.api.routes.get_settings", return_value=cleanup_settings):
        response = client.post("/v1/ops/retention/cleanup")
    assert response.status_code == 200
    payload = response.json()
    assert payload["retention_days"] == 0
    assert "storage_deleted" in payload
