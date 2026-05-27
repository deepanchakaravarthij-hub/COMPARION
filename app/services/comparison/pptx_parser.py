from __future__ import annotations

import hashlib
from dataclasses import dataclass
from io import BytesIO
from typing import Any

from pptx import Presentation  # type: ignore[import-untyped]
from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import-untyped]


@dataclass(frozen=True)
class PptxObject:
    object_id: str
    kind: str
    text: str | None
    table_cells: tuple[tuple[str, ...], ...] | None
    image_sha256: str | None
    image_blob: bytes | None
    shape_signature: str | None
    style_signature: str | None
    bbox: tuple[float, float, float, float] | None

    @property
    def content_signature(self) -> str:
        return (
            f"{self.kind}|{self.text}|{self.table_cells}|{self.image_sha256}|"
            f"{self.shape_signature}|{self.style_signature}"
        )


@dataclass(frozen=True)
class PptxSlide:
    index: int
    title: str | None
    objects: list[PptxObject]

    @property
    def fingerprint(self) -> str:
        payload = "|".join(object_.content_signature for object_ in self.objects)
        digest = hashlib.sha256(payload.encode("utf-8")).hexdigest()
        return f"{self.title}|{len(self.objects)}|{digest}"


@dataclass(frozen=True)
class PptxPresentationModel:
    slides: list[PptxSlide]
    slide_width_emu: int
    slide_height_emu: int


def parse_pptx(content: bytes) -> PptxPresentationModel:
    presentation = Presentation(BytesIO(content))
    slides = [
        _parse_slide(slide, index) for index, slide in enumerate(presentation.slides, start=1)
    ]
    return PptxPresentationModel(
        slides=slides,
        slide_width_emu=int(presentation.slide_width),
        slide_height_emu=int(presentation.slide_height),
    )


def _parse_slide(slide: Any, slide_index: int) -> PptxSlide:
    objects: list[PptxObject] = []
    title_shape = slide.shapes.title
    title = title_shape.text.strip() if title_shape is not None and title_shape.text else None

    for shape in slide.shapes:
        _collect_shape_objects(shape, slide_index, objects, counter=[0])
    return _finalize_slide(slide_index, title, objects)


def _collect_shape_objects(
    shape: Any,
    slide_index: int,
    objects: list[PptxObject],
    *,
    counter: list[int],
    group_prefix: str = "",
) -> None:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for group_index, child_shape in enumerate(shape.shapes, start=1):
            prefix = f"{group_prefix}g{group_index}-"
            _collect_shape_objects(child_shape, slide_index, objects, counter=counter, group_prefix=prefix)
        return

    counter[0] += 1
    object_id = f"s{slide_index}-{group_prefix}o{counter[0]}"
    bbox = _bbox(shape)
    if getattr(shape, "has_text_frame", False):
        text = _normalize_text(shape.text_frame.text if shape.text_frame else "")
        if text:
            objects.append(
                PptxObject(
                    object_id=object_id,
                    kind="text",
                    text=text,
                    table_cells=None,
                    image_sha256=None,
                    image_blob=None,
                    shape_signature=None,
                    style_signature=_text_style_signature(shape),
                    bbox=bbox,
                )
            )
        return

    if getattr(shape, "has_table", False):
        table_cells = _table_cells(shape.table)
        objects.append(
            PptxObject(
                object_id=object_id,
                kind="table",
                text=None,
                table_cells=table_cells,
                image_sha256=None,
                image_blob=None,
                shape_signature=None,
                style_signature=None,
                bbox=bbox,
            )
        )
        return

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        blob = shape.image.blob
        objects.append(
            PptxObject(
                object_id=object_id,
                kind="image",
                text=None,
                table_cells=None,
                image_sha256=hashlib.sha256(blob).hexdigest(),
                image_blob=blob,
                shape_signature=None,
                style_signature=None,
                bbox=bbox,
            )
        )
        return

    objects.append(
        PptxObject(
            object_id=object_id,
            kind="shape",
            text=_normalize_text(getattr(shape, "text", "") or ""),
            table_cells=None,
            image_sha256=None,
            image_blob=None,
            shape_signature=f"{shape.shape_type}|{shape.name}",
            style_signature=None,
            bbox=bbox,
        )
    )


def _finalize_slide(slide_index: int, title: str | None, objects: list[PptxObject]) -> PptxSlide:
    return PptxSlide(index=slide_index, title=title, objects=objects)


def _normalize_text(text: str) -> str:
    return " ".join(text.split())


def _table_cells(table: Any) -> tuple[tuple[str, ...], ...]:
    return tuple(tuple(_normalize_text(cell.text) for cell in row.cells) for row in table.rows)


def _text_style_signature(shape: Any) -> str:
    styles = []
    if not shape.text_frame:
        return ""
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            font = run.font
            styles.append(
                f"{font.name}:{font.size}:{font.bold}:{font.italic}:{font.underline}:"
                f"{_safe_font_color(font)}"
            )
    return "|".join(styles)


def _bbox(shape: Any) -> tuple[float, float, float, float] | None:
    width = getattr(shape, "width", None)
    height = getattr(shape, "height", None)
    left = getattr(shape, "left", None)
    top = getattr(shape, "top", None)
    if width is None or height is None or left is None or top is None:
        return None
    if width == 0 or height == 0:
        return None
    # Values are EMU; normalize to slide coordinates later with a fixed divisor.
    return float(int(left)), float(int(top)), float(int(width)), float(int(height))


def _safe_font_color(font: Any) -> str | None:
    color = getattr(font, "color", None)
    if color is None:
        return None
    try:
        rgb = color.rgb
    except AttributeError:
        return None
    return str(rgb) if rgb else None
